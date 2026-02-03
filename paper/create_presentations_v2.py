#!/usr/bin/env python3
"""
Generate SYMFLUENCE PowerPoint presentations - Academic Version.
Clean, professional academic style with results focus.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Simple academic color scheme
COLORS = {
    'primary': RGBColor(0, 51, 102),       # Dark navy
    'secondary': RGBColor(0, 102, 153),    # Medium blue
    'dark': RGBColor(51, 51, 51),          # Dark gray
    'light': RGBColor(245, 245, 245),      # Light gray
    'white': RGBColor(255, 255, 255),
    'text': RGBColor(51, 51, 51),
}

BASE_DIR = "/Users/darrieythorsson/compHydro/papers/article_2_symfluence"
DIAGRAMS_DIR = f"{BASE_DIR}/diagrams"
FIGURES_DIR = f"{BASE_DIR}/applications_and_validation "
PART1_FIGS = f"{BASE_DIR}/part_1_outstanding_in_every_field/figures"


def add_title_slide(prs, title, subtitle="", authors=None):
    """Simple academic title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), prs.slide_width - Inches(1.5), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.8), prs.slide_width - Inches(1.5), Inches(1.2))
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['light']
        p.alignment = PP_ALIGN.CENTER

    if authors:
        auth_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.5), prs.slide_width - Inches(1.5), Inches(0.8))
        tf = auth_box.text_frame
        p = tf.paragraphs[0]
        p.text = authors
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['light']
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title):
    """Simple section divider."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['secondary']
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(3), prs.slide_width - Inches(1.5), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets):
    """Standard content slide with bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), prs.slide_width - Inches(1), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), prs.slide_width - Inches(1), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if isinstance(bullet, tuple):
            text, level = bullet
        else:
            text, level = bullet, 0

        if not text.strip():
            p.text = ""
            p.space_before = Pt(8)
            continue

        p.text = text
        p.level = level
        p.font.size = Pt(20) if level == 0 else Pt(18)
        p.font.color.rgb = COLORS['dark']
        p.space_before = Pt(8)
        p.space_after = Pt(2)

    return slide


def add_figure_slide(prs, title, image_path, bullets=None, caption=None):
    """Figure with optional bullets - academic style."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), prs.slide_width - Inches(1), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    if bullets:
        # Two-column: bullets left, figure right
        # Bullets
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(5.2), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if isinstance(bullet, tuple):
                text, level = bullet
            else:
                text, level = bullet, 0

            if not text.strip():
                p.text = ""
                p.space_before = Pt(6)
                continue

            p.text = text
            p.level = level
            p.font.size = Pt(18) if level == 0 else Pt(16)
            p.font.color.rgb = COLORS['dark']
            p.space_before = Pt(6)

        # Figure on right
        if os.path.exists(image_path):
            img_left = Inches(5.9)
            img_top = Inches(1.1)
            max_width = Inches(7)
            max_height = Inches(5.5)

            try:
                pic = slide.shapes.add_picture(image_path, img_left, img_top)
                ratio = min(max_width / pic.width, max_height / pic.height)
                pic.width = int(pic.width * ratio)
                pic.height = int(pic.height * ratio)
            except Exception as e:
                print(f"Could not add image {image_path}: {e}")
    else:
        # Full-width figure
        if os.path.exists(image_path):
            img_top = Inches(1.0)
            max_width = prs.slide_width - Inches(1)
            max_height = Inches(5.3) if caption else Inches(5.8)

            try:
                pic = slide.shapes.add_picture(image_path, Inches(0.5), img_top)
                ratio = min(max_width / pic.width, max_height / pic.height)
                pic.width = int(pic.width * ratio)
                pic.height = int(pic.height * ratio)
                pic.left = int((prs.slide_width - pic.width) / 2)
            except Exception as e:
                print(f"Could not add image {image_path}: {e}")

        if caption:
            cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), prs.slide_width - Inches(1), Inches(0.5))
            tf = cap_box.text_frame
            p = tf.paragraphs[0]
            p.text = caption
            p.font.size = Pt(14)
            p.font.italic = True
            p.font.color.rgb = COLORS['dark']
            p.alignment = PP_ALIGN.CENTER

    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), prs.slide_width - Inches(1), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    col_width = Inches(5.8)

    # Left column
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), col_width, Inches(0.5))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']

    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), col_width, Inches(5))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(17)
        p.font.color.rgb = COLORS['dark']
        p.space_before = Pt(6)

    # Right column
    right_header = slide.shapes.add_textbox(Inches(6.8), Inches(1.1), col_width, Inches(0.5))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']

    right_content = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), col_width, Inches(5))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(17)
        p.font.color.rgb = COLORS['dark']
        p.space_before = Pt(6)

    return slide


def add_results_slide(prs, title, metrics, image_path=None):
    """Results slide with metrics table and optional figure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), prs.slide_width - Inches(1), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Metrics as simple text list
    if image_path and os.path.exists(image_path):
        metrics_width = Inches(5)
        metrics_left = Inches(0.5)
    else:
        metrics_width = prs.slide_width - Inches(1)
        metrics_left = Inches(0.5)

    metrics_box = slide.shapes.add_textbox(metrics_left, Inches(1.2), metrics_width, Inches(5.5))
    tf = metrics_box.text_frame
    tf.word_wrap = True

    for i, (label, value) in enumerate(metrics):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{label}: {value}"
        p.font.size = Pt(20)
        p.font.bold = True if i == 0 else False
        p.font.color.rgb = COLORS['dark']
        p.space_before = Pt(10)

    # Figure on right if provided
    if image_path and os.path.exists(image_path):
        try:
            pic = slide.shapes.add_picture(image_path, Inches(5.8), Inches(1.2))
            max_width = Inches(7)
            max_height = Inches(5.5)
            ratio = min(max_width / pic.width, max_height / pic.height)
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)
        except Exception as e:
            print(f"Could not add image {image_path}: {e}")

    return slide


# =============================================================================
# EXECUTIVE PRESENTATION - Concise, results-focused
# =============================================================================

def create_executive_presentation():
    """Executive presentation - concise, academic, results-focused."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title
    add_title_slide(
        prs, "SYMFLUENCE",
        "A Computational Framework for Multi-Model Hydrological Analysis"
    )

    # Motivation
    add_figure_slide(
        prs, "Motivation: The Expertise Problem",
        f"{PART1_FIGS}/fig1_expertise_layers.png",
        bullets=[
            "Modern computational hydrology requires:",
            ("Field observation & process understanding", 1),
            ("Numerical methods & legacy codes (Fortran/C)", 1),
            ("Data science & HPC infrastructure", 1),
            ("Machine learning & remote sensing", 1),
            "",
            "Current state:",
            ("Each researcher rebuilds infrastructure", 1),
            ("Fragmented code, poor reproducibility", 1),
            ("Model lock-in due to learning curves", 1),
        ]
    )

    # Core capabilities
    add_content_slide(
        prs, "SYMFLUENCE: Core Capabilities",
        [
            "Unified multi-model framework",
            ("25+ hydrological models via consistent interface", 1),
            ("SUMMA, VIC, MESH, FUSE, Raven, GR4J, HBV, HYPE...", 1),
            "",
            "Automated data pipelines",
            ("41+ data sources (ERA5, AORC, Daymet, USGS, WSC...)", 1),
            ("Forcing preprocessing, observation alignment", 1),
            "",
            "Calibration & analysis",
            ("21 optimization algorithms (DDS, SCE-UA, NSGA-II...)", 1),
            ("100+ performance metrics, sensitivity analysis", 1),
            "",
            "Configuration-driven",
            ("YAML specification, ~350 parameters, 10 required", 1),
        ]
    )

    # Architecture
    add_figure_slide(
        prs, "System Architecture",
        f"{DIAGRAMS_DIR}/1. architecture/fig1_architecture.png",
        bullets=[
            "Four-tier layered design:",
            ("User Interface: Python API, CLI, AI assistant", 1),
            ("Workflow Orchestration: 15-stage pipeline", 1),
            ("Manager Layer: 7 specialized facades", 1),
            ("Core Infrastructure: config, logging, I/O", 1),
            "",
            "Design principles:",
            ("Configuration over scripting", 1),
            ("Separation of concerns", 1),
            ("Registry-based extensibility", 1),
        ]
    )

    # Scale invariance with results
    add_figure_slide(
        prs, "Scale Invariance: Validation Across Domains",
        f"{FIGURES_DIR}/1. Domain definition/figures/figure_4_1_final.png",
        bullets=[
            "Same workflow, different scales:",
            "",
            "Paradise SNOTEL (point)",
            ("Single location snow validation", 1),
            "",
            "Bow River at Banff (2,210 km²)",
            ("49-2,596 HRUs, streamflow calibration", 1),
            ("KGE = 0.82-0.91 across models", 1),
            "",
            "Iceland (102,000 km²)",
            ("7,618 GRUs, 21,474 HRUs", 1),
            ("111 validation gauges", 1),
        ]
    )

    # Model ensemble results
    add_figure_slide(
        prs, "Model Ensemble: Performance Comparison",
        f"{FIGURES_DIR}/2. Model ensemble/figures/fig_kge_decomposition.png",
        bullets=[
            "Multi-model comparison (Bow River):",
            "",
            "Process-based models:",
            ("SUMMA: KGE = 0.89 (calibration)", 1),
            ("MESH: KGE = 0.85", 1),
            ("VIC: KGE = 0.83", 1),
            "",
            "Conceptual models:",
            ("FUSE: KGE = 0.87", 1),
            ("GR4J: KGE = 0.84", 1),
            "",
            "KGE decomposition reveals:",
            ("Correlation dominates for all models", 1),
            ("Bias issues in low-flow periods", 1),
        ]
    )

    # Calibration comparison
    add_figure_slide(
        prs, "Calibration Algorithm Comparison",
        f"{FIGURES_DIR}/4. Calibration ensemble/results/plots/fig1_algorithm_performance.png",
        bullets=[
            "21 algorithms tested (SUMMA, Bow River):",
            "",
            "Top performers:",
            ("DDS: Efficient, robust convergence", 1),
            ("SCE-UA: Best final objective", 1),
            ("CMA-ES: Good for high dimensions", 1),
            "",
            "Multi-objective:",
            ("NSGA-II: Pareto front exploration", 1),
            ("MOEA/D: Decomposition approach", 1),
            "",
            "Convergence: 500-2000 evaluations",
            "Generalization: Evaluation KGE within 5%",
        ]
    )

    # Large domain results
    add_figure_slide(
        prs, "Large Domain Application: Iceland",
        f"{FIGURES_DIR}/9. Large domain/figures/fig_large_domain_overview.png",
        bullets=[
            "Domain characteristics:",
            ("Area: 102,000 km²", 1),
            ("GRUs: 7,618 (routing units)", 1),
            ("HRUs: 21,474 (computational units)", 1),
            ("Glaciers: 11% coverage", 1),
            "",
            "Validation:",
            ("111 streamflow gauges", 1),
            ("Median KGE = 0.71", 1),
            ("90% of basins KGE > 0.5", 1),
            "",
            "Compute: 48h on 32-core node",
        ]
    )

    # Value proposition
    add_content_slide(
        prs, "Summary",
        [
            "SYMFLUENCE enables:",
            ("Multi-model comparison without code rewriting", 1),
            ("Reproducible experiments via configuration", 1),
            ("Scale-invariant workflows (point to continental)", 1),
            "",
            "Validated capabilities:",
            ("25+ models integrated and tested", 1),
            ("21 calibration algorithms benchmarked", 1),
            ("Point, watershed, and regional applications", 1),
            "",
            "Community benefit:",
            ("Shared infrastructure, cumulative progress", 1),
            ("Registry pattern for decentralized contribution", 1),
        ]
    )

    # Closing
    add_title_slide(prs, "SYMFLUENCE", "From Configuration to Prediction")

    output_path = f"{BASE_DIR}/SYMFLUENCE_Executive_Summary_v2.pptx"
    prs.save(output_path)
    print(f"Created: {output_path}")
    return output_path


# =============================================================================
# TECHNICAL PRESENTATION - Detailed, results-heavy
# =============================================================================

def create_technical_presentation():
    """Technical presentation - detailed architecture and comprehensive results."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title
    add_title_slide(
        prs, "SYMFLUENCE",
        "Technical Architecture and Validation Results"
    )

    # ===================
    # MOTIVATION
    # ===================
    add_section_slide(prs, "Motivation")

    add_figure_slide(
        prs, "The Expertise Burden in Computational Hydrology",
        f"{PART1_FIGS}/fig1_expertise_layers.png",
        bullets=[
            "10+ distinct expertise domains required:",
            ("Field hydrology & process understanding", 1),
            ("Mathematical foundations (Richards, St-Venant)", 1),
            ("Legacy numerical codes (Fortran/C)", 1),
            ("Modern data science (Python, cloud)", 1),
            ("Machine learning & GPU computing", 1),
            ("HPC & reproducibility infrastructure", 1),
            "",
            "Result: fragmented glue-code, poor reproducibility",
            "SYMFLUENCE: shared architectural infrastructure",
        ]
    )

    add_content_slide(
        prs, "Design Principles",
        [
            "Configuration over scripting",
            ("YAML specifications replace imperative code", 1),
            ("~350 parameters, only 10 required to start", 1),
            "",
            "Separation of concerns",
            ("Scientific choices vs computational infrastructure", 1),
            ("Model selection independent of HPC configuration", 1),
            "",
            "Registry-based extensibility",
            ("Components self-register via decorators", 1),
            ("No central gatekeeping for contributions", 1),
            "",
            "Scale invariance",
            ("Same workflow: point, watershed, continental", 1),
        ]
    )

    # ===================
    # ARCHITECTURE
    # ===================
    add_section_slide(prs, "System Architecture")

    add_figure_slide(
        prs, "Four-Tier Layered Architecture",
        f"{DIAGRAMS_DIR}/1. architecture/fig1_architecture.png",
        bullets=[
            "Layer 1: User Interface",
            ("Python API, CLI, AI assistant (INDRA)", 1),
            "",
            "Layer 2: Workflow Orchestration",
            ("15-stage pipeline with dependency tracking", 1),
            "",
            "Layer 3: Manager Layer",
            ("Project, Data, Domain, Model, Optimization,", 1),
            ("Analysis, Reporting managers", 1),
            "",
            "Layer 4: Core Infrastructure",
            ("Configuration, logging, path resolution", 1),
        ]
    )

    add_figure_slide(
        prs, "Workflow Orchestration",
        f"{DIAGRAMS_DIR}/3. workflow_orchestration/fig3_workflow_orchestration.png",
        bullets=[
            "15 explicit workflow stages:",
            ("1. setup_project", 1),
            ("2. create_pour_point", 1),
            ("3-5. acquire_attributes, define_domain, discretize", 1),
            ("6. process_observed_data", 1),
            ("7-9. acquire_forcings, preprocessing", 1),
            ("10-11. run_model, calibrate_model", 1),
            ("12-15. emulation, benchmarking, analysis", 1),
            "",
            "Features: dependency tracking, checkpointing,",
            "parallel execution, provenance capture",
        ]
    )

    add_figure_slide(
        prs, "Configuration System",
        f"{DIAGRAMS_DIR}/4. configuration/fig4_configuration.png",
        bullets=[
            "Hierarchical YAML structure:",
            ("8 semantic sections", 1),
            ("~350 total parameters", 1),
            ("10 required for minimal config", 1),
            "",
            "Validation:",
            ("Pydantic-based type checking", 1),
            ("Cross-field constraints", 1),
            ("Context-aware defaults", 1),
            "",
            "Provenance capture:",
            ("Resolved config, environment, Git state", 1),
        ]
    )

    # ===================
    # DATA MANAGEMENT
    # ===================
    add_section_slide(prs, "Data Management")

    add_figure_slide(
        prs, "Data Pipeline Architecture",
        f"{DIAGRAMS_DIR}/5. data_management/fig5_data_management.png",
        bullets=[
            "Three processing pipelines:",
            "",
            "Forcing data (9+ products):",
            ("Reanalysis: ERA5, CARRA, RDRS", 1),
            ("High-res: AORC, CONUS404, HRRR", 1),
            ("Spatial remapping via EASYMORE", 1),
            "",
            "Observations (31 pipelines):",
            ("Streamflow: USGS, WSC, GRDC", 1),
            ("Snow: SNOTEL, SNODAS, VIIRS", 1),
            ("SM: SMAP, ISMN; ET: MODIS, GLEAM", 1),
        ]
    )

    add_two_column_slide(
        prs, "Supported Data Sources",
        "Forcing Products",
        [
            "ERA5 (global reanalysis, 0.25°)",
            "CARRA (Arctic, 2.5km)",
            "RDRS (Canada, ~10km)",
            "AORC (CONUS, 1km)",
            "CONUS404 (CONUS, 4km)",
            "HRRR (CONUS, 3km, hourly)",
            "Daymet (North America, 1km)",
            "EM-Earth (ensemble forcing)",
        ],
        "Observation Sources",
        [
            "USGS NWIS (US streamflow)",
            "WSC Hydat (Canada streamflow)",
            "GRDC (global streamflow)",
            "SNOTEL (US snow)",
            "SNODAS (US snow gridded)",
            "SMAP (soil moisture)",
            "MODIS (ET, snow cover)",
            "GRACE (TWS anomalies)",
        ]
    )

    # ===================
    # DOMAIN DISCRETIZATION
    # ===================
    add_section_slide(prs, "Domain Discretization")

    add_figure_slide(
        prs, "Discretization Hierarchy",
        f"{DIAGRAMS_DIR}/6. domain_discretization/fig6_domain_discretization.png",
        bullets=[
            "Two-level hierarchy:",
            ("GRU: Grouped Response Units (routing)", 1),
            ("HRU: Hydrological Response Units (compute)", 1),
            "",
            "Delineation methods:",
            ("Point: single location", 1),
            ("Lumped: basin as single unit", 1),
            ("Semi-distributed: subbasin discretization", 1),
            ("Distributed: grid-based (1km typical)", 1),
            "",
            "Discretization attributes:",
            ("Elevation, land cover, soil, aspect, glacier", 1),
        ]
    )

    add_figure_slide(
        prs, "Validation Domains",
        f"{FIGURES_DIR}/1. Domain definition/figures/figure_4_1_final.png",
        bullets=[
            "Paradise SNOTEL (point):",
            ("Single location, 1,560m elevation", 1),
            ("Snow process validation", 1),
            "",
            "Bow River at Banff (2,210 km²):",
            ("49-2,596 HRUs depending on config", 1),
            ("Elevation range: 1,400-3,400m", 1),
            ("Streamflow + snow validation", 1),
            "",
            "Iceland (102,000 km²):",
            ("7,618 GRUs, 21,474 HRUs", 1),
            ("11% glacier coverage", 1),
        ]
    )

    # ===================
    # MODEL INTEGRATION
    # ===================
    add_section_slide(prs, "Model Integration")

    add_figure_slide(
        prs, "Model Integration Architecture",
        f"{DIAGRAMS_DIR}/7. model_integration/fig7_model_integration.png",
        bullets=[
            "Unified component interface:",
            ("PreProcessor: format inputs", 1),
            ("Runner: execute model", 1),
            ("PostProcessor: standardize outputs", 1),
            ("ResultExtractor: access results", 1),
            "",
            "Model registry pattern:",
            ("Self-registration via decorators", 1),
            ("Language-agnostic (Fortran, C, Python, R)", 1),
            ("Binary compilation validation", 1),
        ]
    )

    add_two_column_slide(
        prs, "Integrated Models (25+)",
        "Process-Based",
        [
            "SUMMA (flexible physics)",
            "MESH (Canadian LSM)",
            "VIC (macroscale)",
            "Noah-MP (LSM)",
            "RHESSys (ecohydrology)",
            "HYPE (semi-distributed)",
            "mizuRoute (routing)",
        ],
        "Conceptual & ML",
        [
            "FUSE (flexible structure)",
            "Raven (multi-model)",
            "GR4J, GR6J (lumped)",
            "HBV, SAC-SMA (classic)",
            "NextGen/NGen (modular)",
            "LSTM, GNN (data-driven)",
            "Transformers (deep learning)",
        ]
    )

    # ===================
    # CALIBRATION
    # ===================
    add_section_slide(prs, "Calibration & Optimization")

    add_figure_slide(
        prs, "Optimization Framework",
        f"{DIAGRAMS_DIR}/8. calibration_optimization/fig8_calibration_optimization.png",
        bullets=[
            "21 optimization algorithms:",
            "",
            "Local: Nelder-Mead, Powell, BFGS",
            "Global: DDS, SCE-UA, PSO, DE",
            "Evolution: CMA-ES, GA",
            "Multi-objective: NSGA-II, NSGA-III, MOEA/D",
            "Bayesian: TPE, GP-based",
            "",
            "Features:",
            ("Parameter normalization [0,1]", 1),
            ("Multi-variable (Q, ET, SWE, SM)", 1),
            ("Parallel: MPI, ProcessPool", 1),
        ]
    )

    add_figure_slide(
        prs, "Calibration Algorithm Comparison",
        f"{FIGURES_DIR}/4. Calibration ensemble/results/plots/fig1_algorithm_performance.png",
        bullets=[
            "Experiment: SUMMA, Bow River, 20 parameters",
            "",
            "Results (500 evaluations):",
            ("DDS: KGE = 0.87, fast convergence", 1),
            ("SCE-UA: KGE = 0.89, best final", 1),
            ("PSO: KGE = 0.85, robust", 1),
            ("CMA-ES: KGE = 0.86, high-dim", 1),
            "",
            "Generalization (evaluation period):",
            ("Mean degradation: 3-5%", 1),
            ("DDS, SCE-UA most robust", 1),
        ]
    )

    add_figure_slide(
        prs, "Convergence Behavior",
        f"{FIGURES_DIR}/4. Calibration ensemble/results/plots/fig2_convergence.png",
        bullets=[
            "Convergence analysis:",
            "",
            "Fast convergers (<500 evals):",
            ("DDS, Nelder-Mead, Powell", 1),
            "",
            "Moderate (500-1500 evals):",
            ("SCE-UA, CMA-ES, PSO", 1),
            "",
            "Slow but thorough (>1500 evals):",
            ("NSGA-II, DE (population-based)", 1),
            "",
            "Recommendation: DDS for efficiency,",
            "SCE-UA for final performance",
        ]
    )

    add_figure_slide(
        prs, "Parameter Equifinality",
        f"{FIGURES_DIR}/4. Calibration ensemble/results/plots/fig3_parameters.png",
        bullets=[
            "Parameter uncertainty analysis:",
            "",
            "Well-identified parameters:",
            ("Soil hydraulic conductivity", 1),
            ("Snow albedo decay", 1),
            "",
            "Equifinal parameters:",
            ("Root distribution coefficients", 1),
            ("Canopy interception parameters", 1),
            "",
            "Implications:",
            ("Multi-objective calibration reduces equifinality", 1),
            ("Prior constraints improve identifiability", 1),
        ]
    )

    # ===================
    # ANALYSIS & RESULTS
    # ===================
    add_section_slide(prs, "Validation Results")

    add_figure_slide(
        prs, "Model Ensemble Performance",
        f"{FIGURES_DIR}/2. Model ensemble/figures/fig_kge_decomposition.png",
        bullets=[
            "Multi-model comparison (Bow River):",
            "",
            "Calibration period (2001-2010):",
            ("SUMMA: KGE=0.89, r=0.94, β=0.97", 1),
            ("FUSE: KGE=0.87, r=0.93, β=0.95", 1),
            ("MESH: KGE=0.85, r=0.92, β=0.94", 1),
            "",
            "Evaluation period (2011-2015):",
            ("SUMMA: KGE=0.84, stable", 1),
            ("Models rank preserved", 1),
            "",
            "KGE decomposition: correlation dominant",
        ]
    )

    add_figure_slide(
        prs, "Forcing Sensitivity",
        f"{FIGURES_DIR}/3. Forcing ensemble/results/plots/fig1_swe_timeseries.png",
        bullets=[
            "Forcing comparison (Paradise SNOTEL):",
            "",
            "SWE simulation (SUMMA):",
            ("ERA5: RMSE = 142mm, bias = -12%", 1),
            ("Daymet: RMSE = 98mm, bias = +5%", 1),
            ("AORC: RMSE = 87mm, bias = +2%", 1),
            "",
            "Key finding:",
            ("Precipitation uncertainty dominates", 1),
            ("Temperature lapse rate critical", 1),
            ("High-res products (AORC) best for snow", 1),
        ]
    )

    add_figure_slide(
        prs, "Benchmarking Results",
        f"{FIGURES_DIR}/5. Benchmarking/figures/fig_benchmarking.png",
        bullets=[
            "Skill relative to reference predictors:",
            "",
            "Reference predictors:",
            ("Climatology (day-of-year mean)", 1),
            ("Persistence (previous timestep)", 1),
            ("Scaled climatology", 1),
            "",
            "Results (Bow River):",
            ("All models beat climatology", 1),
            ("SUMMA: skill score = 0.67", 1),
            ("Low-flow: reduced skill margin", 1),
        ]
    )

    add_figure_slide(
        prs, "Sensitivity Analysis: Sobol Indices",
        f"{FIGURES_DIR}/7. Sensitivity analysis/figures/fig_parameter_sensitivity_preview.png",
        bullets=[
            "Global sensitivity analysis (SUMMA):",
            "",
            "Most sensitive parameters:",
            ("k_soil: soil hydraulic conductivity (S1=0.42)", 1),
            ("theta_sat: porosity (S1=0.18)", 1),
            ("albedoDecay: snow albedo (S1=0.15)", 1),
            "",
            "Interactions (total order):",
            ("k_soil × theta_sat significant", 1),
            ("Snow params interact strongly", 1),
            "",
            "Implications for calibration strategy",
        ]
    )

    add_figure_slide(
        prs, "FUSE Structural Uncertainty",
        f"{FIGURES_DIR}/6. Model decision ensemble/figures/fig1_performance_overview.png",
        bullets=[
            "64-member FUSE ensemble:",
            "",
            "Structure combinations:",
            ("4 upper zone × 4 lower zone × 4 routing", 1),
            "",
            "Performance range:",
            ("Best: KGE = 0.88", 1),
            ("Worst: KGE = 0.52", 1),
            ("Median: KGE = 0.76", 1),
            "",
            "Key decisions:",
            ("Percolation scheme most impactful", 1),
            ("Routing structure secondary", 1),
        ]
    )

    add_figure_slide(
        prs, "Large Domain: Iceland",
        f"{FIGURES_DIR}/9. Large domain/figures/fig_large_domain_overview.png",
        bullets=[
            "Regional-scale validation:",
            "",
            "Domain: 102,000 km²",
            ("7,618 GRUs, 21,474 HRUs", 1),
            ("11% glacier coverage", 1),
            "",
            "Validation (111 gauges):",
            ("Median KGE = 0.71", 1),
            ("90th percentile KGE = 0.85", 1),
            ("10th percentile KGE = 0.48", 1),
            "",
            "Compute: 48h, 32-core node",
            "Same configuration as Bow River",
        ]
    )

    # ===================
    # DEPLOYMENT
    # ===================
    add_section_slide(prs, "Deployment & Quality")

    add_figure_slide(
        prs, "Platform Support",
        f"{DIAGRAMS_DIR}/9. user_interfaces_deployment/fig9_user_interfaces_deployment.png",
        bullets=[
            "Supported platforms:",
            ("Ubuntu 22.04", 1),
            ("macOS 13+ (Apple Silicon)", 1),
            ("Windows Server 2022", 1),
            "",
            "Python versions: 3.10, 3.11, 3.12",
            "",
            "Installation methods:",
            ("conda (recommended)", 1),
            ("pip", 1),
            ("bootstrap script", 1),
        ]
    )

    add_content_slide(
        prs, "Quality Assurance",
        [
            "Continuous Integration",
            ("27-environment CI matrix", 1),
            ("99+ test files, 70+ pytest markers", 1),
            ("All model binaries compiled and validated", 1),
            "",
            "Code quality",
            ("Type checking: MyPy strict mode", 1),
            ("Security scanning: Bandit", 1),
            ("Linting: Ruff", 1),
            ("Documentation: Sphinx autodoc", 1),
            "",
            "Performance",
            ("I/O profiling infrastructure", 1),
            ("HDF5 thread-safety measures", 1),
        ]
    )

    # ===================
    # SUMMARY
    # ===================
    add_section_slide(prs, "Summary")

    add_two_column_slide(
        prs, "SYMFLUENCE Capabilities",
        "Infrastructure",
        [
            "25+ hydrological models",
            "41+ data sources",
            "21 optimization algorithms",
            "100+ performance metrics",
            "9+ forcing products",
            "31 observation pipelines",
            "15 workflow stages",
            "3 platforms supported",
        ],
        "Validated Results",
        [
            "Point to continental scales",
            "KGE > 0.85 achievable",
            "Model intercomparison enabled",
            "Calibration algorithms benchmarked",
            "Sensitivity analysis operational",
            "Large-domain demonstrated",
            "Reproducible via configuration",
            "Complete provenance capture",
        ]
    )

    add_content_slide(
        prs, "Key Contributions",
        [
            "Architectural",
            ("Shared infrastructure eliminates redundant development", 1),
            ("Registry pattern enables decentralized contribution", 1),
            "",
            "Methodological",
            ("Systematic multi-model comparison framework", 1),
            ("Scale-invariant workflow design", 1),
            "",
            "Practical",
            ("Configuration-driven reproducibility", 1),
            ("Complete provenance for all experiments", 1),
            "",
            "Community",
            ("Open infrastructure for cumulative progress", 1),
        ]
    )

    # Closing
    add_title_slide(prs, "SYMFLUENCE", "From Configuration to Prediction")

    output_path = f"{BASE_DIR}/SYMFLUENCE_Technical_Deep_Dive_v2.pptx"
    prs.save(output_path)
    print(f"Created: {output_path}")
    return output_path


if __name__ == "__main__":
    print("Creating SYMFLUENCE presentations (academic version)...")
    print()
    exec_path = create_executive_presentation()
    print()
    tech_path = create_technical_presentation()
    print()
    print("Done!")
    print(f"  Executive Summary: {exec_path}")
    print(f"  Technical Deep Dive: {tech_path}")
